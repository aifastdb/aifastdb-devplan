# -*- coding: utf-8 -*-
"""
竞聘PPT自动生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 配色方案
# ============================================================
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)     # 主色-深蓝
COLOR_SECONDARY = RGBColor(0x2B, 0x7A, 0xE8)   # 次色-中蓝
COLOR_ACCENT = RGBColor(0x00, 0xA3, 0xFF)       # 强调色-亮蓝
COLOR_DARK = RGBColor(0x1E, 0x29, 0x3B)         # 深色文字
COLOR_BODY = RGBColor(0x37, 0x47, 0x4F)         # 正文文字
COLOR_LIGHT = RGBColor(0x78, 0x90, 0x9C)        # 浅色文字
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)    # 浅色背景
COLOR_GREEN = RGBColor(0x00, 0xC8, 0x53)        # 绿色-成果
COLOR_ORANGE = RGBColor(0xFF, 0x8F, 0x00)       # 橙色-强调
COLOR_RED = RGBColor(0xE5, 0x39, 0x35)          # 红色-警示
COLOR_GOLD = RGBColor(0xFF, 0xB3, 0x00)         # 金色

# ============================================================
# 字号变量 (修改这里可快速调整全局字号)
# ============================================================
FONT_COVER_TITLE = 44      # 封面大标题
FONT_BIG_NUMBER = 36       # 大数字 / 结语标题
FONT_SECTION_TITLE = 28    # 页面顶部标题栏
FONT_H1 = 22               # 一级标题（目录条目）
FONT_STAT_NUM = 20          # 统计数字 / 圆形序号
FONT_H2 = 18               # 二级标题（板块标题）
FONT_H3 = 16               # 三级标题（卡片标题）
FONT_H4 = 15               # 四级标题（子板块标题）
FONT_SUBTITLE = 14          # 副标题 / 次级文字
FONT_BODY = 13              # 正文
FONT_BODY_SM = 12.5         # 正文（小）
FONT_DETAIL = 12            # 详细内容 / 列表项
FONT_CAPTION = 11           # 标注 / 效果标签
FONT_PAGE_NUM = 10          # 页码

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# 工具函数
# ============================================================

def add_bg_rect(slide, color=COLOR_PRIMARY, left=0, top=0, width=None, height=None):
    """添加背景矩形"""
    w = width or SLIDE_W
    h = height or SLIDE_H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=FONT_H2,
                color=COLOR_BODY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="微软雅黑"):
    """添加文本框并返回text_frame"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=FONT_H3, color=COLOR_BODY,
                  bold=False, space_before=Pt(4), space_after=Pt(2),
                  alignment=PP_ALIGN.LEFT, font_name="微软雅黑", level=0):
    """在text_frame中新增段落"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    p.level = level
    return p


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6),
                   color=COLOR_ACCENT):
    """添加左侧装饰条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide):
    """页面底部装饰条"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()


def add_page_number(slide, num, total):
    """页码"""
    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
                Inches(1), Inches(0.35),
                f"{num}/{total}", font_size=FONT_PAGE_NUM, color=COLOR_LIGHT,
                alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle=""):
    """统一的内容页标题区域"""
    # 顶部白色块
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_WHITE
    bar.line.fill.background()

    # 底部分隔线
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.14), SLIDE_W, Pt(1.5)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    sep.line.fill.background()

    add_textbox(slide, Inches(0.7), Inches(0.15), Inches(10), Inches(0.65),
                title, font_size=FONT_SECTION_TITLE, color=COLOR_DARK, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.72), Inches(10), Inches(0.35),
                    subtitle, font_size=FONT_SUBTITLE, color=COLOR_LIGHT)


def add_card(slide, left, top, width, height, fill_color=COLOR_WHITE):
    """添加卡片矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.5)
    shape.adjustments[0] = 0.04
    return shape


TOTAL_PAGES = 17


# ============================================================
# 第1页：封面
# ============================================================
def make_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, COLOR_PRIMARY)

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.5), Inches(2.85), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(1.5), Inches(3.05), Inches(10), Inches(1.2),
                "竞  聘  汇  报", font_size=FONT_COVER_TITLE, color=COLOR_WHITE, bold=True)

    add_textbox(slide, Inches(1.5), Inches(4.35), Inches(10), Inches(0.7),
                "近三年重点工作业绩  ·  个人荣誉  ·  年度工作推进思路",
                font_size=FONT_H2, color=RGBColor(0xBB, 0xDE, 0xFB))

    # 底部装饰
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08))
    bot.fill.solid()
    bot.fill.fore_color.rgb = COLOR_ACCENT
    bot.line.fill.background()


# ============================================================
# 第2页：目录
# ============================================================
def make_toc():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "汇报提纲")
    add_bottom_bar(slide)
    add_page_number(slide, 2, TOTAL_PAGES)

    items = [
        ("01", "近三年重点工作业绩", "固网指标攻坚 · 重大工程 · 政企支撑 · 安全合规 · 运维管理 · 培训赋能"),
        ("02", "近五年个人荣誉", "荣誉表彰 · 技能认证 · 考核成绩"),
        ("03", "本年度重点工作推进思路", "云化收官 · 国铁割接 · 防火墙换代 · 指标精管 · 政企拓展"),
        ("04", "AI赋能工作展望", "智能运维 · 效率提升 · 未来方向"),
    ]

    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.55) + i * Inches(1.35)
        # 序号圆形
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.08),
                                        Inches(0.65), Inches(0.65))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_PRIMARY
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.paragraphs[0].text = num
        tf_c.paragraphs[0].font.size = Pt(FONT_STAT_NUM)
        tf_c.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf_c.paragraphs[0].font.bold = True
        tf_c.paragraphs[0].font.name = "微软雅黑"
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_textbox(slide, Inches(2.15), y, Inches(9), Inches(0.5),
                    title, font_size=FONT_H1, color=COLOR_DARK, bold=True)
        add_textbox(slide, Inches(2.15), y + Inches(0.5), Inches(9), Inches(0.4),
                    desc, font_size=FONT_BODY, color=COLOR_LIGHT)

        # 分隔线
        if i < 3:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(1.2), y + Inches(1.15), Inches(10.5), Pt(1))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            sep.line.fill.background()


# ============================================================
# 第3页：工作业绩总览
# ============================================================
def make_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "近三年重点工作业绩", "2023年1月 — 2025年12月  ·  核心运维成果总览")
    add_bottom_bar(slide)
    add_page_number(slide, 3, TOTAL_PAGES)

    # 4个关键指标卡片
    cards_data = [
        ("99%", "SCSCF接通率", "从70%提升至99%\n连续24个月达标", COLOR_GREEN),
        ("1,100+", "投诉处理", "三年累计处理\n家集客及铁通投诉", COLOR_PRIMARY),
        ("50+", "呼叫中心管理", "开通/调测/运维\n全生命周期管理", COLOR_SECONDARY),
        ("6项", "重大工程", "国铁割接·云化·SBC升级\n防火墙改造·铁通迁移", COLOR_ORANGE),
    ]

    card_w = Inches(2.6)
    card_h = Inches(2.3)
    start_x = Inches(0.7)
    gap = Inches(0.35)

    for i, (value, title, desc, color) in enumerate(cards_data):
        x = start_x + i * (card_w + gap)
        y = Inches(1.55)
        card = add_card(slide, x, y, card_w, card_h)

        # 顶部色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        add_textbox(slide, x + Inches(0.2), y + Inches(0.25), card_w - Inches(0.4), Inches(0.8),
                    value, font_size=FONT_BIG_NUMBER, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), card_w - Inches(0.4), Inches(0.4),
                    title, font_size=FONT_H3, color=COLOR_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.45), card_w - Inches(0.4), Inches(0.7),
                    desc, font_size=FONT_DETAIL, color=COLOR_LIGHT, alignment=PP_ALIGN.CENTER)

    # 下方核心职责概述
    y2 = Inches(4.2)
    add_accent_bar(slide, Inches(0.7), y2, height=Inches(0.4))
    add_textbox(slide, Inches(1.0), y2, Inches(6), Inches(0.4),
                "核心职责范围", font_size=FONT_H2, color=COLOR_DARK, bold=True)

    items = [
        "IMS固网核心网维护（SCSCF/SBC/ISBC/HSS/ENUM/防火墙）",
        "政企语音专线及呼叫中心全生命周期管理",
        "某热线/短号码/紧急号码等业务接入方案制定与实施",
        "固网指标体系管理（集团对标、质量系统开发、工作台算法）",
        "网络安全合规（等保/漏洞/防火墙策略/反诈）",
        "面向地市技术培训与技能竞赛支撑",
    ]
    tf = add_textbox(slide, Inches(1.0), y2 + Inches(0.5), Inches(11), Inches(2.5),
                     "", font_size=FONT_BODY, color=COLOR_BODY)
    tf.paragraphs[0].text = ""
    for idx, item in enumerate(items):
        p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(FONT_BODY)
        p.font.color.rgb = COLOR_BODY
        p.font.name = "微软雅黑"
        p.space_before = Pt(3)
        p.space_after = Pt(1)


# ============================================================
# 第4页：固网指标攻坚
# ============================================================
def make_indicator():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, '一、固网核心指标攻坚', 'SCSCF接通率从不达标到持续优秀')
    add_bottom_bar(slide)
    add_page_number(slide, 4, TOTAL_PAGES)

    # 左侧：指标变化展示（模拟折线数据点）
    add_accent_bar(slide, Inches(0.7), Inches(1.45), height=Inches(0.4))
    add_textbox(slide, Inches(1.0), Inches(1.45), Inches(5), Inches(0.4),
                "接通率提升历程", font_size=FONT_H2, color=COLOR_DARK, bold=True)

    data_points = [
        ("2023.01", "93.4%", COLOR_ORANGE),
        ("2023.07", "75%", COLOR_RED),
        ("2023.08", "70%（最低）", COLOR_RED),
        ("2023.11", "99%（突破）", COLOR_GREEN),
        ("2024全年", "98~100%", COLOR_GREEN),
        ("2025全年", "99%", COLOR_GREEN),
    ]
    for i, (period, rate, color) in enumerate(data_points):
        y = Inches(2.05) + i * Inches(0.52)
        # 时间
        add_textbox(slide, Inches(1.0), y, Inches(1.6), Inches(0.4),
                    period, font_size=FONT_BODY, color=COLOR_LIGHT, bold=False)
        # 圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.7), y + Inches(0.08),
                                     Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        # 竖线连接
        if i < len(data_points) - 1:
            vl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(2.78), y + Inches(0.26), Pt(2), Inches(0.34))
            vl.fill.solid()
            vl.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            vl.line.fill.background()
        # 数值
        add_textbox(slide, Inches(3.1), y, Inches(2.5), Inches(0.4),
                    rate, font_size=FONT_SUBTITLE, color=color, bold=True)

    # 右侧关键举措
    rx = Inches(6.5)
    add_accent_bar(slide, rx, Inches(1.45), height=Inches(0.4))
    add_textbox(slide, rx + Inches(0.3), Inches(1.45), Inches(5), Inches(0.4),
                "关键举措", font_size=FONT_H2, color=COLOR_DARK, bold=True)

    measures = [
        "建立 数据提取-原因分析-定位IP-联系政企-督促整改 闭环机制",
        "针对智能外呼用户超频外呼/停机仍外呼等问题，反复联系地市处理",
        "深入研究接通率算法，发现CENTREX停机不回复183等技术根因",
        "推动华为U2000升级解决话务量不一致问题",
        "推动核心网工作台固网指标周报开发，实现指标自动化监控",
        "推进质量管理系统ISBC指标报表开发，提供算法、持续验证修正",
    ]

    for i, m in enumerate(measures):
        y = Inches(2.05) + i * Inches(0.65)
        add_card(slide, rx, y, Inches(6.1), Inches(0.55), COLOR_BG_LIGHT)
        add_textbox(slide, rx + Inches(0.15), y + Inches(0.05), Inches(5.8), Inches(0.45),
                    f"✦  {m}", font_size=FONT_BODY_SM, color=COLOR_BODY)

    # 底部成果总结
    result_box = add_card(slide, Inches(0.7), Inches(5.5), Inches(11.9), Inches(0.7),
                          RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, Inches(1.0), Inches(5.55), Inches(11.5), Inches(0.6),
                "✅ 成果：从2023年初多次低于集团80%红线，到2024-2025年连续24个月稳定在99%以上，实现固网核心指标质的飞跃",
                font_size=FONT_SUBTITLE, color=RGBColor(0x2E, 0x7D, 0x32), bold=True)


# ============================================================
# 第5页：重大工程-国铁 & 云化
# ============================================================
def make_projects_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "二、重大工程项目推进（上）", "国铁专网割接改造  ·  IMS固网云化工程")
    add_bottom_bar(slide)
    add_page_number(slide, 5, TOTAL_PAGES)

    # 左侧：国铁
    lx = Inches(0.5)
    add_card(slide, lx, Inches(1.4), Inches(5.9), Inches(5.3))
    add_accent_bar(slide, lx + Inches(0.2), Inches(1.6), height=Inches(0.35), color=COLOR_ORANGE)
    add_textbox(slide, lx + Inches(0.5), Inches(1.58), Inches(5), Inches(0.4),
                "国铁专网割接改造（2024-2025，牵头推进）", font_size=FONT_H3, color=COLOR_DARK, bold=True)

    guotie_items = [
        ("SBC改造", "牵头紧急完成哈尔滨SBC2/5国铁割接功能支持改造，协同承载网完成IP分配，协调多专业方案制定，解决改造中突发问题"),
        ("路由组织", "牵头国铁局数据路由规范研讨、数据制作，组织厂家讨论商定方案，完成关口局/ICSCF局数据制作"),
        ("测试推进", "通过跟踪信令及与集团沟通，解决国铁已割接/未割接号码互通、铁路114短号等拨测问题"),
        ("配套工作", "完成ENUM数据制作、防火墙放通、业支指令验证、95017微信支付割接等"),
    ]
    for i, (title, desc) in enumerate(guotie_items):
        y = Inches(2.15) + i * Inches(1.05)
        add_textbox(slide, lx + Inches(0.5), y, Inches(5.2), Inches(0.3),
                    f"▎{title}", font_size=FONT_BODY, color=COLOR_PRIMARY, bold=True)
        add_textbox(slide, lx + Inches(0.5), y + Inches(0.3), Inches(5.2), Inches(0.65),
                    desc, font_size=FONT_DETAIL, color=COLOR_BODY)

    # 右侧：云化
    rx = Inches(6.8)
    add_card(slide, rx, Inches(1.4), Inches(5.9), Inches(5.3))
    add_accent_bar(slide, rx + Inches(0.2), Inches(1.6), height=Inches(0.35), color=COLOR_SECONDARY)
    add_textbox(slide, rx + Inches(0.5), Inches(1.58), Inches(5), Inches(0.4),
                "IMS固网云化工程（2023-2025，全程参与）", font_size=FONT_H3, color=COLOR_DARK, bold=True)

    yunhua_items = [
        ("HSS云化", "从方案评审→实例化→上线，历时三年；数据提取签保密协议；协调地市遍历性测试；处理12个异常号码"),
        ("网元上云", "完成ICSCF上云、SCSCF上云（大兴安岭等）、ISBG上云，推进固网核心网全面云化"),
        ("业务验证", "协调呼叫中心测试配合、固网游牧数据提供中兴、ENS操作配合，确保云化后业务不中断"),
        ("协同配合", "协调辽宁云科室、网络云、承载网等多专业，完成路由发布、防火墙策略放通等"),
    ]
    for i, (title, desc) in enumerate(yunhua_items):
        y = Inches(2.15) + i * Inches(1.05)
        add_textbox(slide, rx + Inches(0.5), y, Inches(5.2), Inches(0.3),
                    f"▎{title}", font_size=FONT_BODY, color=COLOR_PRIMARY, bold=True)
        add_textbox(slide, rx + Inches(0.5), y + Inches(0.3), Inches(5.2), Inches(0.65),
                    desc, font_size=FONT_DETAIL, color=COLOR_BODY)


# ============================================================
# 第6页：重大工程-SBC升级、防火墙改造、铁通迁移、MGCF割接
# ============================================================
def make_projects_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "二、重大工程项目推进（下）", "SBC升级  ·  防火墙RBM改造  ·  铁通迁移  ·  MGCF业务割接")
    add_bottom_bar(slide)
    add_page_number(slide, 6, TOTAL_PAGES)

    projects = [
        ("全省SBC升级（2023年）", COLOR_PRIMARY,
         ["组织完成全省13个地市、21台SBC设备版本升级",
          "同步完成ISBC软件版本升级及集团工单回复",
          "升级后指标稳定，无业务中断"]),
        ("IMS防火墙RBM改造（2025年）", COLOR_ORANGE,
         ["编写改造方案、计划及发文",
          "完成全省9个地市防火墙RBM改造",
          "完成4A授权更新、综资更新、告警牌更新",
          "编写应急手册，组织9地市应急演练"]),
        ("铁通业务迁移（2023-2024年）", COLOR_SECONDARY,
         ["推进伊春、七台河等地市铁通存量用户迁移",
          "完成佳木斯铁通光缆双路由改造",
          "12地市铁通家客固话下线配合",
          "协调市场部及地市退网事宜"]),
        ("MGCF虚拟呼叫中心割接（2024-2025年）", COLOR_GREEN,
         ["将MGCF承载的虚拟呼叫中心业务全部割接至ISBC",
          "完成1XX09、8321XXX6等号码割接",
          "后续清理中继数据、收回IP地址",
          "400号码路由由MGCF改至ISBC完成"]),
    ]

    card_w = Inches(2.85)
    gap = Inches(0.3)
    for i, (title, color, items) in enumerate(projects):
        x = Inches(0.55) + i * (card_w + gap)
        y = Inches(1.45)
        add_card(slide, x, y, card_w, Inches(5.2))
        # 标题栏
        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.65))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        # 矩形无 adjustments，跳过圆角调整
        add_textbox(slide, x + Inches(0.15), y + Inches(0.08), card_w - Inches(0.3), Inches(0.55),
                    title, font_size=FONT_BODY, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            iy = y + Inches(0.85) + j * Inches(0.65)
            add_textbox(slide, x + Inches(0.15), iy, card_w - Inches(0.3), Inches(0.6),
                        f"▸  {item}", font_size=FONT_DETAIL, color=COLOR_BODY)


# ============================================================
# 第7页：国铁客服专网割接 — 全流程详解
# ============================================================
def make_rail_detail():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "国铁客服专网割接 — 全流程详解", "13地市全省覆盖  ·  端到端闭环管理")
    add_bottom_bar(slide)
    add_page_number(slide, 7, TOTAL_PAGES)

    # 六步流程卡片（2行 x 3列）
    steps = [
        ("01", "割接方案制定", "集团指导方案落地\n专业间对接、方案细化", COLOR_PRIMARY),
        ("02", "业支开户配合", "提供开户指令\n测试验证", COLOR_SECONDARY),
        ("03", "设备配合/资源准备", "SBC选择/功能改造\n新增接入IP、信令分组", COLOR_ACCENT),
        ("04", "割接局数据制作", "ICSCF/BGCF/ENUM\n防火墙/短号码全套数据", COLOR_ORANGE),
        ("05", "链路对接", "ISBC与SIP GW\n物理链路对接方案/数据", COLOR_GREEN),
        ("06", "割接问题处理", "测试链路问题\n各地市割接中问题处理", COLOR_RED),
    ]

    card_w = Inches(3.7)
    card_h = Inches(1.65)
    gap_x = Inches(0.35)
    gap_y = Inches(0.3)

    for i, (num, title, desc, color) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = Inches(0.6) + col * (card_w + gap_x)
        y = Inches(1.45) + row * (card_h + gap_y)

        add_card(slide, x, y, card_w, card_h)
        # 序号圆
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2),
                                        Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.paragraphs[0].text = num
        tf_c.paragraphs[0].font.size = Pt(FONT_H3)
        tf_c.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf_c.paragraphs[0].font.bold = True
        tf_c.paragraphs[0].font.name = "微软雅黑"
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_textbox(slide, x + Inches(0.8), y + Inches(0.15), Inches(2.7), Inches(0.4),
                    title, font_size=FONT_SUBTITLE, color=COLOR_DARK, bold=True)
        add_textbox(slide, x + Inches(0.8), y + Inches(0.55), Inches(2.7), Inches(0.9),
                    desc, font_size=FONT_DETAIL, color=COLOR_BODY)

    # 连接箭头（行内步骤间用小三角）
    for row in range(2):
        for col in range(2):
            ax = Inches(0.6) + (col + 1) * (card_w + gap_x) - gap_x / 2 - Inches(0.1)
            ay = Inches(1.45) + row * (card_h + gap_y) + card_h / 2 - Inches(0.08)
            arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                           ax, ay, Inches(0.2), Inches(0.16))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_LIGHT
            arrow.line.fill.background()
            arrow.rotation = 90.0

    # 底部成果卡片
    results_y = Inches(5.1)
    results = [
        ("全省覆盖", "13个地市", "全部完成割接", COLOR_PRIMARY),
        ("精细管理", "端到端闭环", "方案-数据-测试-上线", COLOR_GREEN),
        ("质量保障", "问题快速响应", "割接过程零重大事故", COLOR_ORANGE),
    ]
    rw = Inches(3.7)
    for i, (r_title, r_val, r_desc, r_color) in enumerate(results):
        rx = Inches(0.6) + i * (rw + Inches(0.35))
        add_card(slide, rx, results_y, rw, Inches(1.3), RGBColor(0xF5, 0xF8, 0xFC))
        # 色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, results_y, rw, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = r_color
        bar.line.fill.background()

        add_textbox(slide, rx + Inches(0.2), results_y + Inches(0.15), rw - Inches(0.4), Inches(0.35),
                    r_title, font_size=FONT_SUBTITLE, color=COLOR_DARK, bold=True)
        add_textbox(slide, rx + Inches(0.2), results_y + Inches(0.5), Inches(1.4), Inches(0.35),
                    r_val, font_size=FONT_STAT_NUM, color=r_color, bold=True)
        add_textbox(slide, rx + Inches(1.7), results_y + Inches(0.55), rw - Inches(1.9), Inches(0.35),
                    r_desc, font_size=FONT_DETAIL, color=COLOR_LIGHT)


# ============================================================
# 第8页：政企业务支撑
# ============================================================
def make_enterprise():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "三、政企业务支撑与收入保障", "呼叫中心管理  ·  某热线接入  ·  行业客户  ·  容量规划")
    add_bottom_bar(slide)
    add_page_number(slide, 8, TOTAL_PAGES)

    # 呼叫中心
    lx = Inches(0.5)
    add_card(slide, lx, Inches(1.4), Inches(6.0), Inches(2.55))
    add_accent_bar(slide, lx + Inches(0.2), Inches(1.55), height=Inches(0.3))
    add_textbox(slide, lx + Inches(0.45), Inches(1.5), Inches(4), Inches(0.35),
                "呼叫中心全生命周期管理", font_size=FONT_H4, color=COLOR_DARK, bold=True)
    cc_tf = add_textbox(slide, lx + Inches(0.45), Inches(1.95), Inches(5.3), Inches(1.8),
                        "", font_size=FONT_DETAIL, color=COLOR_BODY)
    cc_items = [
        "累计完成50+个呼叫中心开通、调测、IP变更、并发速率调整、关闭等全流程",
        "涉及北京嘉信、安徽达策、深圳虎贲、山东弘然等全国多家平台",
        "完成24个95呼叫中心数据清理及MGCF虚拟呼叫中心全部割接",
        "建立白名单管理机制，配合反诈要求完成外呼规范化管理",
    ]
    for item in cc_items:
        p = cc_tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(FONT_BODY_SM)
        p.font.color.rgb = COLOR_BODY
        p.font.name = "微软雅黑"
        p.space_before = Pt(3)

    # 某热线
    rx = Inches(6.8)
    add_card(slide, rx, Inches(1.4), Inches(6.0), Inches(2.55))
    add_accent_bar(slide, rx + Inches(0.2), Inches(1.55), height=Inches(0.3), color=COLOR_ORANGE)
    add_textbox(slide, rx + Inches(0.45), Inches(1.5), Inches(4), Inches(0.35),
                "某热线/短号接入支撑", font_size=FONT_H4, color=COLOR_DARK, bold=True)
    hl_tf = add_textbox(slide, rx + Inches(0.45), Inches(1.95), Inches(5.3), Inches(1.8),
                        "", font_size=FONT_DETAIL, color=COLOR_BODY)
    hl_items = [
        "完成绥化、牡丹江、双鸭山、鹤岗、黑河等10+地市某热线接入方案与实施",
        "实现鹤岗某热线视频功能、大兴安岭1XX33视频呼叫功能",
        "七台河某热线完成9个号码从MGCF割接至ISBC",
        "处理大量某热线外呼/呼入故障（单通、不振铃、联通互通等）",
    ]
    for item in hl_items:
        p = hl_tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(FONT_BODY_SM)
        p.font.color.rgb = COLOR_BODY
        p.font.name = "微软雅黑"
        p.space_before = Pt(3)

    # 行业客户
    y3 = Inches(4.25)
    add_card(slide, lx, y3, Inches(8.6), Inches(2.55))
    add_accent_bar(slide, lx + Inches(0.2), y3 + Inches(0.15), height=Inches(0.3), color=COLOR_GREEN)
    add_textbox(slide, lx + Inches(0.45), y3 + Inches(0.1), Inches(4), Inches(0.35),
                "行业客户专项支撑", font_size=FONT_H4, color=COLOR_DARK, bold=True)

    sectors = [
        ("电力系统", "95598落地码通过ISBC HMR策略实现，SIP GW接入方案制定"),
        ("金融行业", "人保投标技术应答、建行彩铃、农信社400中继方案完成"),
        ("政务系统", "96120割接、公积金呼叫中心(牡丹江/黑河)调测、医保局/税务局接入"),
        ("铁路系统", "国铁专网割接全程支撑，SBC改造、路由组织、测试配合"),
    ]
    for i, (sector, detail) in enumerate(sectors):
        sy = y3 + Inches(0.55) + i * Inches(0.47)
        add_textbox(slide, lx + Inches(0.45), sy, Inches(1.3), Inches(0.35),
                    f"▎{sector}", font_size=FONT_DETAIL, color=COLOR_PRIMARY, bold=True)
        add_textbox(slide, lx + Inches(1.8), sy, Inches(6.5), Inches(0.35),
                    detail, font_size=FONT_BODY_SM, color=COLOR_BODY)

    # ISBC容量
    add_card(slide, Inches(9.4), y3, Inches(3.4), Inches(2.55))
    add_accent_bar(slide, Inches(9.6), y3 + Inches(0.15), height=Inches(0.3), color=COLOR_ACCENT)
    add_textbox(slide, Inches(9.85), y3 + Inches(0.1), Inches(3), Inches(0.35),
                "ISBC容量规划", font_size=FONT_H4, color=COLOR_DARK, bold=True)
    cap_tf = add_textbox(slide, Inches(9.85), y3 + Inches(0.55), Inches(2.8), Inches(1.8),
                         "", font_size=FONT_DETAIL, color=COLOR_BODY)
    for item in ["协同政企、计划共同评估容量需求", "撰写扩容方案材料", "提出软硬件配置建议", "配合设计院提供利用率数据"]:
        p = cap_tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(FONT_DETAIL)
        p.font.color.rgb = COLOR_BODY
        p.font.name = "微软雅黑"
        p.space_before = Pt(3)


# ============================================================
# 第8页：网络安全 + 运维
# ============================================================
def make_security():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "四、网络安全合规 & 投诉处理与日常运维", "安全合规保障  ·  投诉趋势分析  ·  日常运维管理")
    add_bottom_bar(slide)
    add_page_number(slide, 9, TOTAL_PAGES)

    # 左：安全
    lx = Inches(0.5)
    add_card(slide, lx, Inches(1.4), Inches(6.0), Inches(5.3))
    add_accent_bar(slide, lx + Inches(0.2), Inches(1.55), height=Inches(0.35), color=COLOR_RED)
    add_textbox(slide, lx + Inches(0.5), Inches(1.52), Inches(4), Inches(0.35),
                "网络安全与合规保障", font_size=FONT_H3, color=COLOR_DARK, bold=True)

    sec_items = [
        ("漏洞管理", "每月持续华为IMS、华三/华为防火墙漏洞扫描\n中高级漏洞整改及处理报告撰写"),
        ("合规检查", "防火墙策略扫描、合规检查、弱口令核查\n配合集团/管局多次现场检查"),
        ("等保工作", "固网等保系统录入、报告修改、安全资产报表上报\n工信部资产一致性核对"),
        ("反诈支撑", "IMS固网反诈IFC方案制定及实施\n配合反诈签约验证，完成5万+号码加固"),
        ("防火墙管理", "华三防火墙全省升级、策略整改\n应急流程制定、RBM改造、健康检查"),
    ]

    for i, (title, desc) in enumerate(sec_items):
        y = Inches(2.1) + i * Inches(0.92)
        add_textbox(slide, lx + Inches(0.5), y, Inches(1.6), Inches(0.3),
                    f"◆  {title}", font_size=FONT_DETAIL, color=COLOR_PRIMARY, bold=True)
        add_textbox(slide, lx + Inches(2.2), y, Inches(3.8), Inches(0.75),
                    desc, font_size=FONT_DETAIL, color=COLOR_BODY)

    # 右：运维
    rx = Inches(6.8)
    add_card(slide, rx, Inches(1.4), Inches(6.0), Inches(3.0))
    add_accent_bar(slide, rx + Inches(0.2), Inches(1.55), height=Inches(0.35), color=COLOR_SECONDARY)
    add_textbox(slide, rx + Inches(0.5), Inches(1.52), Inches(4), Inches(0.35),
                "投诉处理趋势（逐年下降）", font_size=FONT_H3, color=COLOR_DARK, bold=True)

    # 柱状图模拟
    bar_data = [("2023年", 37, COLOR_ORANGE), ("2024年", 32, COLOR_SECONDARY), ("2025年", 27, COLOR_GREEN)]
    max_val = 40
    bar_area_y = Inches(2.15)
    bar_h_max = Inches(1.6)
    for i, (year, val, color) in enumerate(bar_data):
        bx = rx + Inches(0.8) + i * Inches(1.7)
        ratio = val / max_val
        bh = int(bar_h_max * ratio)
        by = bar_area_y + bar_h_max - bh
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, Inches(1.0), bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, bx, by - Inches(0.3), Inches(1.0), Inches(0.3),
                    f"月均{val}件", font_size=FONT_DETAIL, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, bx, bar_area_y + bar_h_max + Inches(0.05), Inches(1.0), Inches(0.25),
                    year, font_size=FONT_DETAIL, color=COLOR_LIGHT, alignment=PP_ALIGN.CENTER)

    # 日常运维
    add_card(slide, rx, Inches(4.7), Inches(6.0), Inches(2.0))
    add_accent_bar(slide, rx + Inches(0.2), Inches(4.85), height=Inches(0.3), color=COLOR_ACCENT)
    add_textbox(slide, rx + Inches(0.5), Inches(4.82), Inches(4), Inches(0.35),
                "日常运维管理", font_size=FONT_H4, color=COLOR_DARK, bold=True)
    daily_tf = add_textbox(slide, rx + Inches(0.5), Inches(5.25), Inches(5.3), Inches(1.2),
                           "", font_size=FONT_BODY_SM, color=COLOR_BODY)
    for item in [
        "周报月报某热线拨打情况及固网指标部分持续撰写（三年不间断）",
        "配合财务部每月提供全业务话务量数据",
        "与业支月度对账持续跟进（HSS/ENUM号码文件生成）",
        "现场服务工单派单、归档、稽核系统录入、计提全流程管理",
    ]:
        p = daily_tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(FONT_DETAIL)
        p.font.color.rgb = COLOR_BODY
        p.font.name = "微软雅黑"
        p.space_before = Pt(2)


# ============================================================
# 第10页：固话游牧模板变更
# ============================================================
def make_nomadic_template():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "五、固话游牧模板变更 — 全省规模化IP地址整合",
                       "从问题发现到全省存量模板改造  ·  跨部门协同推进")
    add_bottom_bar(slide)
    add_page_number(slide, 10, TOTAL_PAGES)

    # ---- 左侧：推进历程时间线 ----
    lx = Inches(0.5)
    add_card(slide, lx, Inches(1.4), Inches(6.0), Inches(5.3))
    add_accent_bar(slide, lx + Inches(0.2), Inches(1.55), height=Inches(0.35), color=COLOR_SECONDARY)
    add_textbox(slide, lx + Inches(0.5), Inches(1.52), Inches(4), Inches(0.35),
                "工作推进历程", font_size=FONT_H3, color=COLOR_DARK, bold=True)

    timeline_data = [
        ("2025.04", "市场部变更游牧模板，不区分区县，IP地址整合", COLOR_PRIMARY),
        ("2025.06", "地址核对，完成游牧数据制作", COLOR_PRIMARY),
        ("2025.07", "变更脚本数据核查，配合业支程序上线验证", COLOR_SECONDARY),
        ("2025.09", "市场部牵头修改存量模版申请，网络提供指令", COLOR_ORANGE),
        ("2025.10", "固话加固问题及存量游牧模板修改讨论", COLOR_GREEN),
    ]

    for i, (period, desc, color) in enumerate(timeline_data):
        y = Inches(2.1) + i * Inches(0.62)
        # 时间
        add_textbox(slide, lx + Inches(0.4), y, Inches(1.3), Inches(0.35),
                    period, font_size=FONT_DETAIL, color=COLOR_LIGHT, bold=False)
        # 圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx + Inches(1.8), y + Inches(0.06),
                                     Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        # 竖线
        if i < len(timeline_data) - 1:
            vl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        lx + Inches(1.87), y + Inches(0.22), Pt(2), Inches(0.4))
            vl.fill.solid()
            vl.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            vl.line.fill.background()
        # 描述
        add_textbox(slide, lx + Inches(2.15), y, Inches(3.4), Inches(0.5),
                    desc, font_size=FONT_BODY_SM, color=COLOR_BODY)

    # 背景说明
    add_textbox(slide, lx + Inches(0.4), Inches(5.3), Inches(5.2), Inches(1.0),
                "背景：集客号码使用家客IP地址导致注册失败，推动全省游牧模板规范化改造",
                font_size=FONT_DETAIL, color=COLOR_LIGHT)

    # ---- 右侧：关键举措 + 成果 ----
    rx = Inches(6.8)
    add_card(slide, rx, Inches(1.4), Inches(6.0), Inches(3.1))
    add_accent_bar(slide, rx + Inches(0.2), Inches(1.55), height=Inches(0.35), color=COLOR_ORANGE)
    add_textbox(slide, rx + Inches(0.5), Inches(1.52), Inches(4), Inches(0.35),
                "关键举措", font_size=FONT_H3, color=COLOR_DARK, bold=True)

    measures = [
        "根因分析：定位家客用户动态获取临区IP导致注册失败",
        "跨部门协同：协调市场部、政企部、业支、地市多方联动",
        "IP地址整合：不区分区县，按地市维度统一整合IP地址段",
        "自动化脚本：开发游牧模板变更脚本，批量数据制作与核查",
        "存量改造：推动全省存量游牧模板修改纳入正式流程",
    ]
    for i, m in enumerate(measures):
        y = Inches(2.05) + i * Inches(0.43)
        add_textbox(slide, rx + Inches(0.4), y, Inches(5.4), Inches(0.4),
                    f"✦  {m}", font_size=FONT_DETAIL, color=COLOR_BODY)

    # 成果卡片
    add_card(slide, rx, Inches(4.75), Inches(6.0), Inches(1.95), RGBColor(0xE8, 0xF5, 0xE9))
    add_accent_bar(slide, rx + Inches(0.2), Inches(4.9), height=Inches(0.3), color=COLOR_GREEN)
    add_textbox(slide, rx + Inches(0.5), Inches(4.87), Inches(4), Inches(0.35),
                "核心成果", font_size=FONT_H4, color=COLOR_DARK, bold=True)
    results_tf = add_textbox(slide, rx + Inches(0.5), Inches(5.3), Inches(5.3), Inches(1.2),
                             "", font_size=FONT_DETAIL, color=RGBColor(0x2E, 0x7D, 0x32))
    for item in [
        "解决家客用户获取临区IP导致游牧限制无法注册投诉",
        "完成全省多地市游牧模板变更数据制作",
        "实现游牧模板变更脚本化、自动化",
        "推动存量游牧模板修改纳入市场部牵头的正式流程",
    ]:
        p = results_tf.add_paragraph()
        p.text = f"✅  {item}"
        p.font.size = Pt(FONT_BODY_SM)
        p.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
        p.font.name = "微软雅黑"
        p.space_before = Pt(3)


# ============================================================
# 第11页：培训
# ============================================================
def make_training():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "六、培训授课与团队建设", "知识沉淀与传承")
    add_bottom_bar(slide)
    add_page_number(slide, 11, TOTAL_PAGES)

    trainings = [
        ("2024.03", "VONR语音业务L3培训授课", "科室", COLOR_PRIMARY),
        ("2024.04", "VoLTE信令流程L2课程讲解", "科室", COLOR_PRIMARY),
        ("2024.10", "《IMS固网业务实训》面向地市授课", "全省地市", COLOR_SECONDARY),
        ("2025.04", "L2、L3实训课程授课", "全省地市", COLOR_SECONDARY),
        ("2025.07", "华北大区比武监考裁判", "大区级", COLOR_ORANGE),
        ("2025.08", "新员工授课", "新入职员工", COLOR_GREEN),
    ]

    for i, (date, title, scope, color) in enumerate(trainings):
        y = Inches(1.55) + i * Inches(0.9)
        add_card(slide, Inches(1.2), y, Inches(10.9), Inches(0.75))

        # 时间标签
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.5), y + Inches(0.15), Inches(1.3), Inches(0.42))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tf_tag = tag.text_frame
        tf_tag.paragraphs[0].text = date
        tf_tag.paragraphs[0].font.size = Pt(FONT_DETAIL)
        tf_tag.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf_tag.paragraphs[0].font.bold = True
        tf_tag.paragraphs[0].font.name = "微软雅黑"
        tf_tag.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_textbox(slide, Inches(3.1), y + Inches(0.15), Inches(6.5), Inches(0.45),
                    title, font_size=FONT_H3, color=COLOR_DARK, bold=True)

        # 范围标签
        scope_tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(9.8), y + Inches(0.15), Inches(2.0), Inches(0.42))
        scope_tag.fill.solid()
        scope_tag.fill.fore_color.rgb = COLOR_BG_LIGHT
        scope_tag.line.fill.background()
        tf_s = scope_tag.text_frame
        tf_s.paragraphs[0].text = scope
        tf_s.paragraphs[0].font.size = Pt(FONT_DETAIL)
        tf_s.paragraphs[0].font.color.rgb = COLOR_LIGHT
        tf_s.paragraphs[0].font.name = "微软雅黑"
        tf_s.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_s.vertical_anchor = MSO_ANCHOR.MIDDLE


# ============================================================
# 第12页：个人荣誉（占位）
# ============================================================
def make_honors():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "近五年个人荣誉", "2021年 — 2025年  ·  荣誉表彰与技能认证")
    add_bottom_bar(slide)
    add_page_number(slide, 12, TOTAL_PAGES)

    add_textbox(slide, Inches(2.5), Inches(3.0), Inches(8), Inches(1.5),
                "【请在此处补充近五年个人荣誉内容】\n\n如：年度考核结果、荣誉称号、技能竞赛获奖、专利/论文/创新成果、集团/省公司表彰等",
                font_size=FONT_H2, color=COLOR_LIGHT, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第13-14页：年度推进思路
# ============================================================
def make_plan_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "本年度重点工作推进思路（一）", "固网云化收官  ·  国铁专网全面推进")
    add_bottom_bar(slide)
    add_page_number(slide, 13, TOTAL_PAGES)

    # 云化收官
    lx = Inches(0.5)
    add_card(slide, lx, Inches(1.4), Inches(6.0), Inches(5.3))
    hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, Inches(1.4), Inches(6.0), Inches(0.55))
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = COLOR_PRIMARY
    hbar.line.fill.background()
    add_textbox(slide, lx + Inches(0.2), Inches(1.45), Inches(5.5), Inches(0.45),
                "一、固网云化收官与稳定运行保障", font_size=FONT_H3, color=COLOR_WHITE, bold=True)

    plans1 = [
        ("目标", "确保IMS固网云化工程按期高质量完成，保障云化后网络稳定运行"),
        ("HSS云化收尾", "完成全业务验证，处理存量异常数据，确保业支双发后所有场景正常"),
        ("固网上云扫尾", "推进剩余网元上云操作，协调承载网、云基础等多专业配合"),
        ("运维体系建立", "建立云化环境下的巡检标准、告警规则、应急预案"),
        ("性能基线建立", "上云后建立新的性能基线，确保指标不低于物理机时代"),
    ]
    for i, (title, desc) in enumerate(plans1):
        y = Inches(2.15) + i * Inches(0.88)
        if i == 0:
            add_textbox(slide, lx + Inches(0.3), y, Inches(5.4), Inches(0.6),
                        f"🎯 {desc}", font_size=FONT_DETAIL, color=COLOR_PRIMARY, bold=True)
        else:
            add_textbox(slide, lx + Inches(0.3), y, Inches(1.8), Inches(0.3),
                        f"▎{title}", font_size=FONT_DETAIL, color=COLOR_PRIMARY, bold=True)
            add_textbox(slide, lx + Inches(0.3), y + Inches(0.3), Inches(5.4), Inches(0.5),
                        desc, font_size=FONT_BODY_SM, color=COLOR_BODY)

    # 国铁
    rx = Inches(6.8)
    add_card(slide, rx, Inches(1.4), Inches(6.0), Inches(5.3))
    hbar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, Inches(1.4), Inches(6.0), Inches(0.55))
    hbar2.fill.solid()
    hbar2.fill.fore_color.rgb = COLOR_ORANGE
    hbar2.line.fill.background()
    add_textbox(slide, rx + Inches(0.2), Inches(1.45), Inches(5.5), Inches(0.45),
                "二、国铁专网割接全面推进", font_size=FONT_H3, color=COLOR_WHITE, bold=True)

    plans2 = [
        ("目标", "完成国铁专网割接全省推广，实现铁路固话业务平稳迁移"),
        ("经验复制", "将哈尔滨SBC改造经验形成标准化方案，推广至其他地市"),
        ("局数据规范化", "基于已积累的路由组织经验，建立省内国铁局数据制作标准流程"),
        ("分批割接", "按照集团部署，制定分批割接计划，每批次方案→测试→割接→验证闭环"),
        ("风险管控", "建立割接回退机制，确保业务零中断"),
    ]
    for i, (title, desc) in enumerate(plans2):
        y = Inches(2.15) + i * Inches(0.88)
        if i == 0:
            add_textbox(slide, rx + Inches(0.3), y, Inches(5.4), Inches(0.6),
                        f"🎯 {desc}", font_size=FONT_DETAIL, color=COLOR_ORANGE, bold=True)
        else:
            add_textbox(slide, rx + Inches(0.3), y, Inches(1.8), Inches(0.3),
                        f"▎{title}", font_size=FONT_DETAIL, color=COLOR_ORANGE, bold=True)
            add_textbox(slide, rx + Inches(0.3), y + Inches(0.3), Inches(5.4), Inches(0.5),
                        desc, font_size=FONT_BODY_SM, color=COLOR_BODY)


def make_plan_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "本年度重点工作推进思路（二）", "防火墙换代  ·  指标精管  ·  政企拓展")
    add_bottom_bar(slide)
    add_page_number(slide, 14, TOTAL_PAGES)

    plans = [
        ("三、防火墙换代与安全能力提升", COLOR_RED,
         "完成48台出保防火墙的替换/续保，提升安全防护水平",
         [("预算落实", "推进26年华三防火墙维保预算审批落地"),
          ("改造巩固", "持续监控已改造防火墙运行状态，完善应急手册"),
          ("策略优化", "定期策略清理、合规检查，减少无效策略"),
          ("安全联动", "加强与安全室协同，建立固网安全事件快速响应机制")]),
        ("四、固网指标精细化管理", COLOR_GREEN,
         'SCSCF接通率持续保持99%以上，从被动处理向主动预防转型',
         [("监控前移", "利用工作台和质量系统，建立日级指标监测预警"),
          ("根因常态化", "建立接通率波动自动分析机制，快速定位影响源"),
          ("地市协同", "完善与政企/地市的协同流程，缩短处理时长"),
          ("存量清理", "持续推进销户仍注册用户清理、游牧模板规范化")]),
        ("五、政企业务支撑能力提升", COLOR_SECONDARY,
         "提升支撑效率和方案质量，助力收入增长",
         [("方案模板化", "某热线接入、呼叫中心开通等高频方案标准化"),
          ("ISBC拓展", "推进容量扩展评估，满足3万并发等大需求"),
          ("COS接入", "配合在线公司COS平台接入ISBC"),
          ("培训赋能", "持续开展面向地市的IMS固网业务培训")]),
    ]

    card_w = Inches(3.85)
    for i, (title, color, goal, items) in enumerate(plans):
        x = Inches(0.4) + i * (card_w + Inches(0.25))
        y = Inches(1.4)
        add_card(slide, x, y, card_w, Inches(5.3))

        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.55))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        add_textbox(slide, x + Inches(0.15), y + Inches(0.07), card_w - Inches(0.3), Inches(0.45),
                    title, font_size=FONT_BODY, color=COLOR_WHITE, bold=True)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.6),
                    f"🎯 {goal}", font_size=FONT_DETAIL, color=color, bold=True)

        for j, (sub_title, sub_desc) in enumerate(items):
            sy = y + Inches(1.5) + j * Inches(0.9)
            add_textbox(slide, x + Inches(0.2), sy, card_w - Inches(0.4), Inches(0.3),
                        f"▎{sub_title}", font_size=FONT_DETAIL, color=color, bold=True)
            add_textbox(slide, x + Inches(0.2), sy + Inches(0.3), card_w - Inches(0.4), Inches(0.5),
                        sub_desc, font_size=FONT_DETAIL, color=COLOR_BODY)


# ============================================================
# 第15页：AI赋能（现有）
# ============================================================
def make_ai_current():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "AI赋能工作展望（一）", "现有工作中可结合AI的场景")
    add_bottom_bar(slide)
    add_page_number(slide, 15, TOTAL_PAGES)

    scenarios = [
        ("投诉故障诊断", "基于信令日志AI自动分析\n构建故障知识库，辅助定位根因", "故障定位时间\n从小时级→分钟级", COLOR_PRIMARY),
        ("接通率指标分析", "AI自动识别接通率异常模式\n自动关联失败原因TOP-N", "从月底发现\n→实时预警", COLOR_SECONDARY),
        ("周报/月报撰写", "大模型自动汇总指标和话务量\n生成报告初稿", "每月节省\n4-6小时", COLOR_GREEN),
        ("话单/信令分析", "AI批量分析异常话单模式\n自动识别盗打/超频等行为", "异常发现效率\n提升10倍+", COLOR_ORANGE),
        ("工单脚本制作", "AI根据需求自动生成局数据脚本\n人工审核后下发", "减少制作时间\n降低人工错误", COLOR_ACCENT),
        ("防火墙策略审计", "AI自动识别冗余/违规/过期策略", "安全合规检查\n效率大幅提升", COLOR_RED),
    ]

    for i, (title, desc, effect, color) in enumerate(scenarios):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.5) + row * Inches(2.8)
        card_w = Inches(3.9)
        card_h = Inches(2.5)

        add_card(slide, x, y, card_w, card_h)
        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.5))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = color
        hbar.line.fill.background()
        add_textbox(slide, x + Inches(0.15), y + Inches(0.07), card_w - Inches(0.3), Inches(0.4),
                    f"💡  {title}", font_size=FONT_SUBTITLE, color=COLOR_WHITE, bold=True)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.3), Inches(0.8),
                    desc, font_size=FONT_DETAIL, color=COLOR_BODY)

        # 效果标签
        ebox = add_card(slide, x + Inches(2.5), y + Inches(0.65), Inches(1.25), Inches(0.7),
                        RGBColor(0xE8, 0xF5, 0xE9))
        add_textbox(slide, x + Inches(2.5), y + Inches(0.7), Inches(1.25), Inches(0.6),
                    effect, font_size=FONT_CAPTION, color=COLOR_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第16页：AI赋能（未来）
# ============================================================
def make_ai_future():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "AI赋能工作展望（二）", "未来工作中AI融合方向")
    add_bottom_bar(slide)
    add_page_number(slide, 16, TOTAL_PAGES)

    futures = [
        ("智能运维 AIOps", "基于历史告警和故障数据，建立预测性维护模型\n在故障发生前预警，降低故障率", "⬆ 网络可靠性"),
        ("智能知识库", "将三年1100+件投诉处理经验结构化\n构建AI可检索的运维知识库", "⬆ 经验传承"),
        ("智能调测助手", "AI辅助呼叫中心接入调测，自动比对信令参数\n提示缺少PAI参数/编解码不匹配等常见错误", "⬆ 调测效率"),
        ("自动化对账", "AI自动完成与业支、财务的月度数据对账\n自动识别差异项，替代手动对账", "⬇ 人力投入"),
        ("智能容量预测", "基于历史话务量趋势，AI预测ISBC/SBC容量需求\n提前规划扩容，避免容量瓶颈", "⬆ 资源效率"),
        ("AI辅助方案生成", "输入客户需求，AI匹配历史方案模板\n自动生成接入方案初稿", "⬆ 方案效率"),
    ]

    for i, (title, desc, tag) in enumerate(futures):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(1.5) + row * Inches(2.8)
        card_w = Inches(3.9)
        card_h = Inches(2.5)

        add_card(slide, x, y, card_w, card_h)

        # 顶部渐变效果（用色条代替）
        gradient_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06))
        gradient_bar.fill.solid()
        gradient_bar.fill.fore_color.rgb = COLOR_ACCENT
        gradient_bar.line.fill.background()

        add_textbox(slide, x + Inches(0.15), y + Inches(0.2), card_w - Inches(0.3), Inches(0.35),
                    f"🚀  {title}", font_size=FONT_SUBTITLE, color=COLOR_DARK, bold=True)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.65), card_w - Inches(0.3), Inches(1.2),
                    desc, font_size=FONT_BODY_SM, color=COLOR_BODY)

        # 效果标签
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           x + Inches(0.15), y + Inches(1.95),
                                           Inches(1.5), Inches(0.35))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
        tag_shape.line.fill.background()
        tf_tag = tag_shape.text_frame
        tf_tag.paragraphs[0].text = tag
        tf_tag.paragraphs[0].font.size = Pt(FONT_DETAIL)
        tf_tag.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        tf_tag.paragraphs[0].font.bold = True
        tf_tag.paragraphs[0].font.name = "微软雅黑"
        tf_tag.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE


# ============================================================
# 第15页：结语
# ============================================================
def make_ending():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, COLOR_PRIMARY)

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.5), Inches(2.5), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(1.5), Inches(2.75), Inches(10), Inches(1.0),
                "扎根技术  ·  深耕运维  ·  拥抱创新", font_size=FONT_BIG_NUMBER, color=COLOR_WHITE, bold=True)

    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
                "以三年实践为基石，以创新思维为引领\n持续为IMS固网高质量运维贡献力量",
                font_size=FONT_H2, color=RGBColor(0xBB, 0xDE, 0xFB))

    add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.6),
                "感谢聆听", font_size=FONT_SECTION_TITLE, color=COLOR_WHITE, bold=True)

    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08))
    bot.fill.solid()
    bot.fill.fore_color.rgb = COLOR_ACCENT
    bot.line.fill.background()


# ============================================================
# 生成PPT
# ============================================================
if __name__ == "__main__":
    make_cover()              # 1
    make_toc()                # 2
    make_overview()           # 3
    make_indicator()          # 4
    make_projects_1()         # 5
    make_projects_2()         # 6
    make_rail_detail()        # 7  (新增：国铁割接全流程详解)
    make_enterprise()         # 8
    make_security()           # 9
    make_nomadic_template()   # 10 (新增：游牧模板变更)
    make_training()           # 11
    make_honors()             # 12
    make_plan_1()             # 13
    make_plan_2()             # 14
    make_ai_current()         # 15
    make_ai_future()          # 16
    make_ending()             # 17

    # 自动版本递增，防止文件被占用时报 PermissionError
    base_dir = r"D:\Project\git"
    base_name = "竞聘汇报PPT"
    version = 1
    while True:
        suffix = f"_v{version}" if version > 1 else ""
        output_path = os.path.join(base_dir, f"{base_name}{suffix}.pptx")
        try:
            prs.save(output_path)
            break
        except PermissionError:
            version += 1
            if version > 20:
                print("❌ 无法保存，请关闭已打开的PPT文件后重试")
                break

    print(f"✅ PPT 已生成: {output_path}")
    print(f"   共 {len(prs.slides)} 页")
