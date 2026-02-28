# -*- coding: utf-8 -*-
"""
竞聘汇报PPT自动生成脚本 (v4 - 现代卡片风格，无顶部横线)
17 页完整竞聘汇报
设计语言：左侧色条 + 柔和阴影 + 干净白卡片 + 无任何顶部色条横线
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 配色
# ============================================================
C_PRI      = RGBColor(0x1A, 0x56, 0xDB)   # 主蓝
C_PRI_D    = RGBColor(0x0F, 0x3D, 0x9E)   # 深蓝
C_PRI_L    = RGBColor(0x3B, 0x82, 0xF6)   # 亮蓝
C_ACCENT   = RGBColor(0x00, 0xA3, 0xFF)   # 强调蓝

C_GREEN    = RGBColor(0x10, 0xB9, 0x81)   # 绿
C_GREEN_D  = RGBColor(0x05, 0x9E, 0x6F)   # 深绿
C_ORANGE   = RGBColor(0xF5, 0x9E, 0x0B)   # 橙
C_RED      = RGBColor(0xEF, 0x44, 0x44)   # 红
C_PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)   # 紫
C_TEAL     = RGBColor(0x14, 0xB8, 0xA6)   # 蓝绿
C_INDIGO   = RGBColor(0x63, 0x66, 0xF1)   # 靛蓝

C_DARK     = RGBColor(0x1E, 0x29, 0x3B)   # 标题
C_BODY     = RGBColor(0x37, 0x47, 0x4F)   # 正文
C_MUTED    = RGBColor(0x64, 0x74, 0x8B)   # 弱化
C_LIGHT    = RGBColor(0x94, 0xA3, 0xB8)   # 浅色
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BG       = RGBColor(0xF8, 0xFA, 0xFC)   # 页面背景
C_SHADOW   = RGBColor(0xD1, 0xD5, 0xDB)   # 阴影
C_CARD     = RGBColor(0xFF, 0xFF, 0xFF)   # 卡片

SW = Inches(13.333)
SH = Inches(7.5)
TOTAL = 17

# ============================================================
# 基础工具
# ============================================================

def set_bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def _sh(s, st, l, t, w, h, fill=None, line=None, lw=None):
    sh = s.shapes.add_shape(st, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        if lw: sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def rect(s, l, t, w, h, fill=None): return _sh(s, MSO_SHAPE.RECTANGLE, l, t, w, h, fill)
def rrect(s, l, t, w, h, fill=None, line=None): return _sh(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill, line)
def oval(s, l, t, w, h, fill=None): return _sh(s, MSO_SHAPE.OVAL, l, t, w, h, fill)

def txt(s, l, t, w, h, text, sz=14, c=C_BODY, b=False, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.alignment = al
    return tb

def blist(s, l, t, w, h, items, sz=11, c=C_BODY, sp=6):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "● " + item
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.space_before = Pt(sp)
        p.space_after = Pt(sp)
    return tb

def pgn(s, n):
    txt(s, Inches(12.2), Inches(7.0), Inches(0.9), Inches(0.25),
        f"{n}/{TOTAL}", sz=9, c=C_LIGHT, al=PP_ALIGN.RIGHT)

def bottom(s):
    rect(s, Inches(0), Inches(7.2), SW, Inches(0.3), fill=C_PRI)

def num_icon(s, x, y, n, c, size=0.42):
    sz = Inches(size)
    o = oval(s, x, y, sz, sz, fill=c)
    p = o.text_frame.paragraphs[0]
    p.text = str(n)
    p.font.size = Pt(int(size * 28))
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# ============================================================
# 现代卡片组件 — 核心设计语言
# ============================================================

def card(s, x, y, w, h, accent=None):
    """现代卡片：阴影 + 白圆角 + 左侧色条（无顶部横线）"""
    # 阴影层
    rrect(s, x + Inches(0.03), y + Inches(0.03), w, h, fill=C_SHADOW)
    # 白卡片
    rrect(s, x, y, w, h, fill=C_CARD)
    # 左侧色条（圆角矩形内部，紧贴左边）
    if accent:
        rect(s, x + Inches(0.01), y + Inches(0.12), Inches(0.055), h - Inches(0.24), fill=accent)


def data_card(s, x, y, w, h, label, value, desc, vc):
    """数据卡片"""
    card(s, x, y, w, h, vc)
    txt(s, x + Inches(0.25), y + Inches(0.18), w - Inches(0.5), Inches(0.25),
        label, sz=10, c=C_MUTED)
    txt(s, x + Inches(0.25), y + Inches(0.42), w - Inches(0.5), Inches(0.5),
        value, sz=28, c=vc, b=True)
    if desc:
        txt(s, x + Inches(0.25), y + Inches(0.95), w - Inches(0.5), Inches(0.25),
            desc, sz=9, c=C_LIGHT)


def titled_card(s, x, y, w, h, title, items, accent):
    """带标题的内容卡片（左色条 + 标题 + bullet列表）"""
    card(s, x, y, w, h, accent)
    # 标题区：圆形小点 + 标题
    oval(s, x + Inches(0.2), y + Inches(0.2), Inches(0.12), Inches(0.12), fill=accent)
    txt(s, x + Inches(0.42), y + Inches(0.14), w - Inches(0.7), Inches(0.3),
        title, sz=15, c=C_DARK, b=True)
    blist(s, x + Inches(0.22), y + Inches(0.55), w - Inches(0.44), h - Inches(0.7),
          items, sz=11, c=C_BODY)


def panel_card(s, x, y, w, h, title, items, accent):
    """面板卡片（左侧宽色带 + 标题 + 列表）"""
    # 阴影
    rrect(s, x + Inches(0.03), y + Inches(0.03), w, h, fill=C_SHADOW)
    # 白卡片
    rrect(s, x, y, w, h, fill=C_CARD)
    # 左侧宽色带
    rect(s, x, y + Inches(0.08), Inches(0.55), h - Inches(0.16), fill=accent)
    # 色带里放图标文字
    txt(s, x + Inches(0.02), y + Inches(0.3), Inches(0.52), Inches(0.3),
        "◆", sz=16, c=C_WHITE, al=PP_ALIGN.CENTER)
    # 标题
    txt(s, x + Inches(0.72), y + Inches(0.2), w - Inches(1), Inches(0.35),
        title, sz=16, c=C_DARK, b=True)
    # 列表
    blist(s, x + Inches(0.72), y + Inches(0.65), w - Inches(1), h - Inches(0.8),
          items, sz=11, c=C_BODY)


def header(s, title, sub=None):
    """页面标题"""
    # 左侧装饰块
    rect(s, Inches(0.55), Inches(0.45), Inches(0.08), Inches(0.55), fill=C_PRI)
    txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
        title, sz=26, c=C_DARK, b=True)
    if sub:
        txt(s, Inches(0.8), Inches(0.92), Inches(10), Inches(0.3),
            sub, sz=13, c=C_MUTED)


# ============================================================
# 17页
# ============================================================

def slide_01(prs):
    """封面"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_PRI)

    # 顶部装饰线
    rect(s, Inches(0), Inches(0), SW, Inches(0.05), fill=C_ACCENT)

    # 中间白色区域（圆角）
    rrect(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(3.5), fill=C_WHITE)
    # 左侧色条装饰
    rect(s, Inches(2.5), Inches(2.2), Inches(0.07), Inches(3.1), fill=C_ACCENT)

    txt(s, Inches(3), Inches(2.5), Inches(7.3), Inches(1),
        "竞  聘  汇  报", sz=44, c=C_PRI, b=True, al=PP_ALIGN.CENTER)
    txt(s, Inches(3), Inches(3.7), Inches(7.3), Inches(0.6),
        "近三年重点工作业绩  |  个人荣誉  |  年度推进思路",
        sz=18, c=C_MUTED, al=PP_ALIGN.CENTER)

    rect(s, Inches(0), Inches(7.0), SW, Inches(0.5), fill=C_PRI_D)


def slide_02(prs):
    """提纲"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)

    txt(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
        "汇 报 提 纲", sz=32, c=C_DARK, b=True)
    rect(s, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.04), fill=C_PRI)

    items = [
        ("01", "近三年重点工作业绩", "固网核心指标 | 重大工程 | 政企支撑 | 安全合规", C_PRI_L),
        ("02", "近五年个人荣誉", "年度考核 | 获奖表彰", C_GREEN),
        ("03", "本年度重点工作推进思路", "云化收官 | 国铁推进 | 指标精管 | 政企拓展", C_ORANGE),
        ("04", "AI赋能展望", "现有工作AI结合 | 未来AI融合方向", C_PURPLE),
    ]

    for i, (num, title, desc, clr) in enumerate(items):
        y = Inches(1.7) + Inches(i * 1.3)
        card(s, Inches(1.3), y, Inches(10.2), Inches(1.0), clr)
        num_icon(s, Inches(1.6), y + Inches(0.22), num, clr, 0.52)
        txt(s, Inches(2.35), y + Inches(0.15), Inches(7), Inches(0.4),
            title, sz=20, c=C_DARK, b=True)
        txt(s, Inches(2.35), y + Inches(0.58), Inches(7), Inches(0.3),
            desc, sz=12, c=C_MUTED)

    bottom(s)
    pgn(s, 2)


def slide_03(prs):
    """核心成果总览"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG)
    header(s, "近三年重点工作业绩", "2023 - 2025 核心成果总览")

    cards_data = [
        ("固网SCSCF接通率", "99%+", "从70%提升至持续优秀", C_GREEN),
        ("投诉处理量", "1100+件", "三年累计 月均降至27件", C_PRI_L),
        ("呼叫中心接入", "50+家", "政企业务累计支撑", C_TEAL),
        ("重大工程", "6项", "国铁/云化/SBC升级等", C_ORANGE),
    ]
    cw, ch, sx, gap = Inches(2.85), Inches(1.35), Inches(0.6), Inches(0.23)
    for i, (lb, v, d, vc) in enumerate(cards_data):
        data_card(s, sx + (cw + gap) * i, Inches(1.55), cw, ch, lb, v, d, vc)

    blocks = [
        ("核心指标攻坚", [
            "SCSCF接通率从70%提升至99%+",
            "建立闭环排障机制",
            "与集团指标核对、话务量校准",
        ], C_PRI_L),
        ("重大工程实施", [
            "国铁专网全省割接（13地市）",
            "IMS固网云化迁移",
            "全省SBC升级及防火墙改造",
        ], C_GREEN),
        ("政企业务支撑", [
            "50+呼叫中心接入管理",
            "某热线全省14地市接入",
            "400/95号码等行业客户支撑",
        ], C_ORANGE),
    ]
    bw = Inches(3.75)
    for i, (title, items, clr) in enumerate(blocks):
        x = Inches(0.6) + (bw + Inches(0.25)) * i
        titled_card(s, x, Inches(3.45), bw, Inches(3.35), title, items, clr)

    bottom(s)
    pgn(s, 3)


def slide_04(prs):
    """固网核心指标"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "一、固网核心指标攻坚", "SCSCF接通率从不达标到持续优秀")

    # 左侧柱状图卡片
    card(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.3))
    txt(s, Inches(0.8), Inches(1.65), Inches(5), Inches(0.35),
        "SCSCF接通率提升历程", sz=14, c=C_DARK, b=True)

    bar_data = [
        ("2023.01", 93, C_ORANGE), ("2023.07", 75, C_RED),
        ("2023.08", 70, C_RED), ("2023.11", 99, C_GREEN),
        ("2024H1", 98, C_GREEN), ("2024H2", 100, C_GREEN),
        ("2025", 99, C_GREEN),
    ]
    bw, bg, by, mh, sx = Inches(0.62), Inches(0.1), Inches(5.6), Inches(3.2), Inches(0.85)
    for i, (label, pct, clr) in enumerate(bar_data):
        x = sx + (bw + bg) * i
        h = mh * (pct / 100)
        y = by - h
        rrect(s, x, y, bw, h, fill=clr)
        txt(s, x - Inches(0.05), y - Inches(0.28), bw + Inches(0.1), Inches(0.28),
            f"{pct}%", sz=10, c=clr, b=True, al=PP_ALIGN.CENTER)
        txt(s, x - Inches(0.1), by + Inches(0.06), bw + Inches(0.2), Inches(0.2),
            label, sz=8, c=C_LIGHT, al=PP_ALIGN.CENTER)

    # 达标线
    ly = by - mh * 0.8
    rect(s, sx - Inches(0.08), ly, Inches(5.3), Inches(0.015), fill=C_RED)
    txt(s, sx + Inches(4.85), ly - Inches(0.12), Inches(1), Inches(0.25),
        "80%达标线", sz=8, c=C_RED, b=True)

    # 右侧举措
    card(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3), C_PRI_L)
    txt(s, Inches(7.3), Inches(1.65), Inches(5), Inches(0.35),
        "关键举措", sz=14, c=C_DARK, b=True)

    measures = [
        "建立 数据提取→原因分析→定位IP→联系政企→督促整改 闭环机制",
        "针对智能外呼用户未开并发/停机仍外呼，联系各地市处理",
        "深入研究接通率算法，发现CENTREX停机不回复183等技术根因",
        "与集团核对话务量不一致问题，推动华为U2000升级解决",
        "推动地市未使用号码摘机处理，从源头降低无效呼叫",
        "2024年起接通率稳定98%-100%，远超集团80%达标线",
    ]
    blist(s, Inches(7.3), Inches(2.1), Inches(5.2), Inches(4.2), measures, sz=11, c=C_BODY, sp=10)

    bottom(s)
    pgn(s, 4)


def slide_05(prs):
    """重大工程（上）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "二、重大工程实施（上）", "国铁专网割接 + IMS固网云化")

    panel_card(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5),
               "国铁专网割接（13地市全省覆盖）", [
                   "完成全省13个地市国铁客服专网割接方案制定与实施",
                   "割接方案：集团指导方案落地，专业间对接、方案细化",
                   "业支开户：提供开户指令，测试验证",
                   "设备配合：SBC选择、功能改造/新增接入IP实现信令分组",
                   "局数据制作：ICSCF/BGCF/ENUM/防火墙/短号码全套",
                   "链路对接：ISBC与SIP GW物理链路对接方案",
                   "持续跟进割接问题处理，确保各地市平稳过渡",
                   "方案细化到IP地址段分配、信令分组等精细化管理",
               ], C_PRI)

    panel_card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.5),
               "IMS固网云化迁移", [
                   "ISBC虚拟化方案审核与优化",
                   "确认CMNET接入点方案，统一CE接入架构",
                   "华为CS/HSS系统账号交接与退网管理",
                   "爱立信IMS设备退网下电全流程执行",
                   "云侧故障处理（医保局/民意调查局中继中断等）",
                   "新旧平台并行期业务连续性保障",
                   "2025年持续推进VM替代及容量规划",
                   "配合云化后承载网改造与链路调整",
               ], C_GREEN)

    bottom(s)
    pgn(s, 5)


def slide_06(prs):
    """重大工程（下）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "二、重大工程实施（下）", "SBC升级 + 防火墙改造 + 铁通迁移 + MGCF割接")

    quads = [
        ("全省SBC升级", C_PRI_L, [
            "完成全省13个地市SBC逐步升级",
            "升级过程中处理SBC注册/呼叫异常",
            "配合SBC容灾方案与DNS方案评估",
        ]),
        ("防火墙RBM改造", C_ORANGE, [
            "华三防火墙升级与漏洞扫描",
            "管局迎检准备与策略清理",
            "2025年推进防火墙批量换代招标",
        ]),
        ("铁通用户迁移", C_GREEN, [
            "铁通专网割接方案持续推进",
            "铁通存量用户迁移支撑",
            "低效地市业务量核查与优化",
        ]),
        ("MGCF割接", C_RED, [
            "MGCF异局址割接至新VM",
            "割接过程通话中断问题攻关",
            "爱立信设备退网全流程管理",
        ]),
    ]
    for i, (t, c, items) in enumerate(quads):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.5) + row * Inches(2.85)
        titled_card(s, x, y, Inches(5.95), Inches(2.55), t, items, c)

    bottom(s)
    pgn(s, 6)


def slide_07(prs):
    """国铁割接全流程"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "国铁客服专网割接 — 全流程详解", "13地市全省覆盖  |  端到端闭环管理")

    steps = [
        ("① 割接方案\n   制定", ["集团指导方案落地", "专业间对接", "方案细化"]),
        ("② 业支开户\n   配合", ["提供开户指令", "测试验证"]),
        ("③ 设备配合\n   资源准备", ["SBC选择/功能改造", "新增接入IP", "信令分组"]),
        ("④ 割接局数据\n   制作", ["ICSCF/BGCF/ENUM", "防火墙/短号码", "全套数据"]),
        ("⑤ 链路对接", ["ISBC与SIP GW", "物理链路", "对接方案/数据"]),
        ("⑥ 割接问题\n   处理", ["测试链路问题", "各地市割接中", "问题处理"]),
    ]

    blues = [
        RGBColor(0x1E, 0x40, 0x8A), RGBColor(0x1E, 0x50, 0x9A),
        RGBColor(0x1E, 0x60, 0xAA), RGBColor(0x2B, 0x70, 0xBB),
        RGBColor(0x38, 0x82, 0xCC), RGBColor(0x45, 0x94, 0xDD),
    ]

    fw, fh = Inches(1.72), Inches(0.65)
    aw, gap = Inches(0.3), Inches(0.12)
    sx, fy, dy = Inches(0.45), Inches(1.55), Inches(2.50)

    for i, (st, sd) in enumerate(steps):
        x = sx + i * (fw + aw + gap)
        # 步骤框（圆角+颜色填充）
        sh = _sh(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, fy, fw, fh, fill=blues[i])
        sh.text_frame.word_wrap = True
        p = sh.text_frame.paragraphs[0]
        p.text = st
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            _sh(s, MSO_SHAPE.RIGHT_ARROW, x + fw + gap * 0.4, fy + Inches(0.18),
                aw, Inches(0.22), fill=C_ACCENT)

        # 详情卡片
        card(s, x, dy, fw, Inches(1.4))
        # 左侧细色条
        rect(s, x + Inches(0.01), dy + Inches(0.12), Inches(0.04), Inches(1.16), fill=blues[i])
        blist(s, x + Inches(0.12), dy + Inches(0.08), fw - Inches(0.2), Inches(1.2),
              sd, sz=9, c=C_BODY, sp=3)

    # 成果区
    sep_y = Inches(4.35)
    txt(s, Inches(0.6), sep_y + Inches(0.1), Inches(4), Inches(0.35),
        "核心成果与亮点", sz=16, c=C_DARK, b=True)

    results = [
        ("全省覆盖", "13个地市", "全部完成割接",
         "覆盖哈尔滨、齐齐哈尔、牡丹江等全部地市", C_PRI_L, "13"),
        ("精细管理", "端到端闭环", "方案→数据→测试→上线",
         "全流程把控每个环节质量", C_GREEN, "E2E"),
        ("质量保障", "零重大事故", "问题快速响应",
         "提前预案、快速响应、业务连续", C_ORANGE, "0"),
    ]
    rcw, rch = Inches(3.75), Inches(2.35)
    for i, (t, v, d, ex, clr, ic) in enumerate(results):
        rx = Inches(0.5) + i * (rcw + Inches(0.25))
        ry = sep_y + Inches(0.55)
        card(s, rx, ry, rcw, rch, clr)
        num_icon(s, rx + Inches(0.2), ry + Inches(0.2), ic, clr, 0.48)
        txt(s, rx + Inches(0.85), ry + Inches(0.18), rcw - Inches(1.1), Inches(0.3),
            t, sz=14, c=clr, b=True)
        txt(s, rx + Inches(0.85), ry + Inches(0.5), rcw - Inches(1.1), Inches(0.3),
            v, sz=20, c=C_DARK, b=True)
        txt(s, rx + Inches(0.85), ry + Inches(0.85), rcw - Inches(1.1), Inches(0.25),
            d, sz=11, c=C_MUTED)
        txt(s, rx + Inches(0.2), ry + Inches(1.3), rcw - Inches(0.4), Inches(0.7),
            ex, sz=9, c=C_LIGHT)

    bottom(s)
    pgn(s, 7)


def slide_08(prs):
    """政企业务支撑"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "三、政企业务支撑", "呼叫中心管理 | 某热线接入 | 行业客户支撑")

    quads = [
        ("呼叫中心管理", Inches(0.5), Inches(1.5), C_PRI_L, [
            "50+家呼叫中心接入运维管理",
            "北京嘉信/安徽达策/湖南数据等并发量核查",
            "呼叫中心直连话单计费验证",
            "400电话/95号码/1XX63方案测试",
            "涉诈号码管控与外呼权限管理",
        ]),
        ("某热线接入", Inches(6.8), Inches(1.5), C_TEAL, [
            "全省14个地市某热线接入与运维",
            "处理单通/外呼失败等复杂技术问题",
            "牡丹江视频接入需求技术评估",
            "某热线日报/周报/月报持续输出",
            "平台侧SIP协议问题（PAI/rtcp参数）攻关",
        ]),
        ("行业客户支撑", Inches(0.5), Inches(4.2), C_GREEN, [
            "人民银行1XX63方案测试与话单验证",
            "人保/二发电等集团用户问题处理",
            "医保局/民意调查局中继中断恢复",
            "智路科技等外呼失败问题排查",
            "政企固话云视讯会议需求支撑",
        ]),
        ("容量规划与管理", Inches(6.8), Inches(4.2), C_ORANGE, [
            "SBC资源分配与Session能力评估",
            "VM资源申请替代物理SBC推进",
            "每月话务量统计分析（4G/5G/固话）",
            "固话一致性核对（每月10日全量比对）",
            "低效无效资产清理上报",
        ]),
    ]
    for t, x, y, clr, items in quads:
        titled_card(s, x, y, Inches(5.95), Inches(2.5), t, items, clr)

    bottom(s)
    pgn(s, 8)


def slide_09(prs):
    """安全合规与投诉"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "四、网络安全合规与投诉处理", "安全体系建设 + 投诉量持续下降")

    panel_card(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5),
               "网络安全合规", [
                   "华为IMS/CS/HSS/防火墙定期漏洞扫描",
                   "IMS等保调整，固网部分优化",
                   "防火墙管局迎检与策略清理（清除无用策略8条+）",
                   "非涉敏系统接入表格填写（CS/IMS/HSS）",
                   "涉诈号码管控：关闭外呼权限、集团双跨问题协调",
                   "反诈评估材料：国际来话/网间来话算法加密整改",
                   "骚扰电话拦截平台核心网对接",
                   "信安部安全需求持续响应",
               ], C_PRI)

    # 投诉处理
    rx = Inches(6.8)
    card(s, rx, Inches(1.5), Inches(5.9), Inches(5.5), C_GREEN)
    # 左侧色带
    rect(s, rx, Inches(1.58), Inches(0.55), Inches(5.34), fill=C_GREEN)
    txt(s, rx + Inches(0.02), Inches(1.8), Inches(0.52), Inches(0.3),
        "◆", sz=16, c=C_WHITE, al=PP_ALIGN.CENTER)
    txt(s, rx + Inches(0.72), Inches(1.65), Inches(5), Inches(0.35),
        "投诉处理趋势", sz=16, c=C_DARK, b=True)
    txt(s, rx + Inches(0.72), Inches(2.15), Inches(5), Inches(0.3),
        "月均投诉量（件/月）", sz=12, c=C_DARK, b=True)

    complaint = [("2023", 37, C_RED), ("2024", 32, C_ORANGE), ("2025", 27, C_GREEN)]
    cbw, cbx, cby, cmh = Inches(1.2), rx + Inches(1.0), Inches(5.0), Inches(2.0)

    for i, (label, val, clr) in enumerate(complaint):
        x = cbx + i * (cbw + Inches(0.25))
        h = cmh * (val / 40)
        y = cby - h
        rrect(s, x, y, cbw, h, fill=clr)
        txt(s, x, y - Inches(0.28), cbw, Inches(0.28),
            f"{val}件", sz=14, c=clr, b=True, al=PP_ALIGN.CENTER)
        txt(s, x, cby + Inches(0.06), cbw, Inches(0.22),
            label + "年", sz=11, c=C_BODY, al=PP_ALIGN.CENTER)

    txt(s, rx + Inches(0.72), Inches(5.35), Inches(5), Inches(0.7),
        "三年投诉持续下降，月均从37件降至27件  ↓27%\n投诉处理经验沉淀，解决效率显著提升",
        sz=11, c=C_BODY)

    bottom(s)
    pgn(s, 9)


def slide_10(prs):
    """游牧模板变更"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "五、固话游牧模板变更", "从问题发现到全省存量模板改造 | 跨部门协同推进")

    # 左侧时间线
    card(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5))
    txt(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
        "固话游牧模板控制用户注册IP范围。2024年初发现\n企宽开集客业务时，集客号码使用家客IP导致注册失败。",
        sz=10, c=C_MUTED)
    txt(s, Inches(0.8), Inches(2.2), Inches(5), Inches(0.3),
        "工作推进历程", sz=14, c=C_PRI_L, b=True)

    timeline = [
        ("2025.04", "市场部变更游牧模板配合，IP地址整合", "规模变更"),
        ("2025.06", "地址核对，完成游牧数据制作", "持续推进"),
        ("2025.07", "游牧模板变更脚本数据核查，业支上线验证", "自动化上线"),
        ("2025.09", "市场部牵头重新提修改存量模版申请", "存量改造"),
        ("2025.10", "固话加固及存量游牧模板修改讨论", "深化完善"),
    ]

    lx = Inches(1.4)
    sy = Inches(2.65)
    sh = Inches(0.7)
    rect(s, lx + Inches(0.03), sy, Inches(0.02), sh * len(timeline) - Inches(0.15), fill=C_PRI_L)

    for i, (date, desc, phase) in enumerate(timeline):
        y = sy + i * sh
        oval(s, lx - Inches(0.04), y - Inches(0.06), Inches(0.14), Inches(0.14), fill=C_PRI_L)
        txt(s, Inches(0.5), y - Inches(0.08), Inches(0.85), Inches(0.22),
            date, sz=9, c=C_PRI_L, b=True, al=PP_ALIGN.RIGHT)
        txt(s, lx + Inches(0.22), y - Inches(0.1), Inches(3.2), Inches(0.45),
            desc, sz=9, c=C_DARK)
        # 标签
        rrect(s, lx + Inches(3.5), y - Inches(0.04), Inches(0.85), Inches(0.22),
              fill=RGBColor(0xEB, 0xF5, 0xFF))
        txt(s, lx + Inches(3.52), y - Inches(0.04), Inches(0.82), Inches(0.22),
            phase, sz=7, c=C_PRI_L, al=PP_ALIGN.CENTER)

    # 右侧
    rx = Inches(6.8)
    card(s, rx, Inches(1.5), Inches(5.9), Inches(5.5), C_PRI_L)
    txt(s, rx + Inches(0.22), Inches(1.65), Inches(5), Inches(0.3),
        "关键举措", sz=14, c=C_PRI_L, b=True)

    measures = [
        "根因分析：定位家客用户动态获取临区IP导致注册失败",
        "跨部门协同：协调市场部、政企部、业支、地市网络部联动",
        "IP地址整合：不区分区县，按地市维度统一整合IP地址段",
        "自动化脚本：开发游牧模板变更脚本，批量数据制作核查",
        "业支程序配合：推动业支开发游牧模板开户功能自动化",
        "存量改造：推动全省存量游牧模板修改纳入正式流程",
    ]
    blist(s, rx + Inches(0.22), Inches(2.0), Inches(5.3), Inches(2.5), measures, sz=11, c=C_BODY, sp=6)

    txt(s, rx + Inches(0.22), Inches(4.6), Inches(5), Inches(0.3),
        "核心成果", sz=14, c=C_PRI_L, b=True)

    results = [
        ("历史遗留解决", "解决家客用户获取临区IP\n导致游牧限制无法注册投诉"),
        ("全省覆盖", "完成多地市游牧模板\n变更数据制作"),
        ("自动化实现", "游牧模板变更\n脚本化、自动化"),
        ("流程规范化", "存量游牧模板修改\n纳入市场部正式流程"),
    ]
    for i, (rt, rd) in enumerate(results):
        col, row = i % 2, i // 2
        cx = rx + Inches(0.15) + col * Inches(2.85)
        cy = Inches(5.0) + row * Inches(1.15)
        cw2, ch2 = Inches(2.65), Inches(1.0)
        card(s, cx, cy, cw2, ch2)
        txt(s, cx + Inches(0.15), cy + Inches(0.08), cw2 - Inches(0.3), Inches(0.22),
            rt, sz=11, c=C_PRI_L, b=True)
        txt(s, cx + Inches(0.15), cy + Inches(0.32), cw2 - Inches(0.3), Inches(0.55),
            rd, sz=9, c=C_BODY)

    bottom(s)
    pgn(s, 10)


def slide_11(prs):
    """培训授课"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "六、培训授课与团队建设", "知识沉淀与传承")

    trainings = [
        ("2023.07", "IMS固网用户开户及问题排查", "面向全省运维人员"),
        ("2024.04", "华为CS核心网工具使用培训", "新交付厂家工程师"),
        ("2024.06", "RBM防火墙测评实操培训", "防火墙专项安全"),
        ("2024.08", "IMS固网架构及接入流程", "面向政企/家客技术"),
        ("2024.10", "接通率指标分析方法论", "IMS运维团队"),
        ("2025.04", "国铁割接流程培训", "面向地市维护团队"),
    ]

    line_y = Inches(3.35)
    rect(s, Inches(0.8), line_y, Inches(11.5), Inches(0.04), fill=C_PRI_L)

    for i, (date, title, target) in enumerate(trainings):
        x = Inches(1.0) + i * Inches(1.95)
        is_top = (i % 2 == 0)
        # 节点
        oval(s, x + Inches(0.65), line_y - Inches(0.09), Inches(0.22), Inches(0.22), fill=C_PRI_L)
        card_y = line_y - Inches(1.85) if is_top else line_y + Inches(0.35)
        cw2, ch2 = Inches(1.75), Inches(1.65)
        card(s, x, card_y, cw2, ch2, C_PRI_L)
        txt(s, x + Inches(0.12), card_y + Inches(0.1), cw2 - Inches(0.24), Inches(0.22),
            date, sz=10, c=C_PRI_L, b=True)
        txt(s, x + Inches(0.12), card_y + Inches(0.38), cw2 - Inches(0.24), Inches(0.65),
            title, sz=10, c=C_DARK, b=True)
        txt(s, x + Inches(0.12), card_y + Inches(1.1), cw2 - Inches(0.24), Inches(0.4),
            target, sz=9, c=C_MUTED)

    # 总结
    card(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.55))
    txt(s, Inches(1.5), Inches(6.25), Inches(10.3), Inches(0.45),
        "累计授课6次  |  覆盖全省运维/政企/家客技术团队  |  推动知识体系化沉淀",
        sz=14, c=C_PRI_L, b=True, al=PP_ALIGN.CENTER)

    bottom(s)
    pgn(s, 11)


def slide_12(prs):
    """个人荣誉"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "近五年个人荣誉", "年度考核 | 获奖表彰")

    card(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5))
    txt(s, Inches(2), Inches(3.5), Inches(9.3), Inches(1),
        "请在此处补充近五年个人荣誉", sz=24, c=C_LIGHT, al=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(4.5), Inches(9.3), Inches(0.8),
        "建议内容：年度考核评级、技术比武获奖、优秀员工、先进个人等",
        sz=14, c=C_LIGHT, al=PP_ALIGN.CENTER)

    bottom(s)
    pgn(s, 12)


def slide_13(prs):
    """年度推进（一）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "本年度重点工作推进思路（一）", "云化收官 + 国铁全面推进")

    panel_card(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5),
               "IMS固网云化收官", [
                   "推进剩余物理SBC VM替代方案落地",
                   "完成爱立信旧设备退网清理",
                   "云侧容量评估与资源动态调整",
                   "新VM稳定性监控与性能基线建立",
                   "承载网配套改造同步推进",
                   "建立云化后运维SOP手册",
               ], C_PRI)

    panel_card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.5),
               "国铁专网全面推进", [
                   "完成剩余地市国铁割接（目标全省100%）",
                   "割接后质量回访与稳定性跟踪",
                   "建立标准化割接SOP模板",
                   "问题知识库沉淀，提高后续效率",
                   "新增地市快速复制已有割接经验",
                   "配合国铁新业务需求评估与接入",
               ], C_GREEN)

    bottom(s)
    pgn(s, 13)


def slide_14(prs):
    """年度推进（二）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_WHITE)
    header(s, "本年度重点工作推进思路（二）", "防火墙换代 + 指标精管 + 政企拓展")

    quads = [
        ("防火墙批量换代", C_ORANGE, [
            "推进防火墙批量换代招标",
            "新旧设备平滑切换方案",
            "安全策略迁移与优化",
        ]),
        ("指标精细化管理", C_PRI_L, [
            "SCSCF接通率持续保持99%以上",
            "从被动处理向主动预防转型",
            "自动化告警与智能预判机制建设",
        ]),
        ("政企业务拓展", C_GREEN, [
            "呼叫中心规模化接入能力提升",
            "新业务场景技术预研",
            "客户自助排障工具开发",
        ]),
        ("团队能力建设", C_TEAL, [
            "持续开展技术培训",
            "运维知识库体系化建设",
            "新人快速上手指南编写",
        ]),
    ]
    for i, (t, c, items) in enumerate(quads):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.5) + row * Inches(2.85)
        titled_card(s, x, y, Inches(5.95), Inches(2.55), t, items, c)

    bottom(s)
    pgn(s, 14)


def _six_grid(s, scenarios, pg):
    """六宫格 — 现代卡片（左色条+圆标+内容）"""
    cw, ch = Inches(3.85), Inches(2.55)
    gx, gy = Inches(0.22), Inches(0.28)
    for i, (title, desc, clr) in enumerate(scenarios):
        col, row = i % 3, i // 3
        x = Inches(0.5) + col * (cw + gx)
        y = Inches(1.5) + row * (ch + gy)

        card(s, x, y, cw, ch, clr)
        num_icon(s, x + Inches(0.2), y + Inches(0.2), i + 1, clr, 0.40)
        txt(s, x + Inches(0.72), y + Inches(0.22), Inches(2.8), Inches(0.3),
            title, sz=14, c=C_DARK, b=True)

        lines = desc.split("\n")
        blist(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), ch - Inches(0.85),
              lines, sz=11, c=C_BODY, sp=5)

    bottom(s)
    pgn(s, pg)


def slide_15(prs):
    """AI赋能（一）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG)
    header(s, "AI赋能展望（一）", "现有工作中可结合AI的场景")

    scenarios = [
        ("投诉工单智能分析",
         "AI自动分类投诉类型，提取关键信息\n推荐解决方案，缩短处理时间50%+", C_PRI_L),
        ("接通率智能预警",
         "基于历史数据训练模型\n提前预判接通率异常趋势，变被动为主动", C_GREEN),
        ("割接方案智能审核",
         "AI辅助审核局数据一致性\n自动检测配置冲突，降低人工遗漏风险", C_ORANGE),
        ("话务量智能预测",
         "ML模型预测话务量趋势\n辅助容量规划与资源动态调配", C_TEAL),
        ("安全漏洞自动扫描",
         "AI增强漏洞扫描，自动关联威胁情报\n生成修复优先级", C_RED),
        ("日报周报自动生成",
         "AI自动采集数据，生成标准化报表\n释放重复性劳动", C_PURPLE),
    ]
    _six_grid(s, scenarios, 15)


def slide_16(prs):
    """AI赋能（二）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG)
    header(s, "AI赋能展望（二）", "未来工作中AI融合方向")

    futures = [
        ("智能运维AIOps",
         "构建IMS网络智能运维平台\n实现故障自动发现-定位-修复闭环", C_PRI_L),
        ("数字孪生网络",
         "建立IMS固网数字孪生模型\n仿真验证割接方案可行性", C_GREEN),
        ("智能客服助手",
         "AI赋能某热线，实现智能应答\n工单自动派发与进度追踪", C_ORANGE),
        ("知识图谱构建",
         "将运维经验结构化为知识图谱\n实现经验智能检索与传承", C_TEAL),
        ("预测性维护",
         "设备健康度AI评估，提前预判\n硬件故障，零停机更换", C_RED),
        ("自然语言运维",
         "语音/文字指令直接操作网管\n降低运维门槛，提升效率", C_PURPLE),
    ]
    _six_grid(s, futures, 16)


def slide_17(prs):
    """结语"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_PRI)
    rect(s, Inches(0), Inches(0), SW, Inches(0.05), fill=C_ACCENT)

    # 中央白色卡片
    rrect(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(3.5), fill=C_WHITE)
    # 左侧色条
    rect(s, Inches(2.5), Inches(2.2), Inches(0.07), Inches(3.1), fill=C_ACCENT)

    txt(s, Inches(3), Inches(2.5), Inches(7.3), Inches(0.8),
        "扎根技术  深耕运维  拥抱创新",
        sz=36, c=C_PRI, b=True, al=PP_ALIGN.CENTER)
    txt(s, Inches(3), Inches(3.5), Inches(7.3), Inches(1.3),
        "以三年的技术积累为基础，以持续创新为动力\n"
        "全力保障固网语音业务高质量运行\n"
        "用AI赋能传统运维，推动数智化转型",
        sz=16, c=C_BODY, al=PP_ALIGN.CENTER)

    txt(s, Inches(2), Inches(6.1), Inches(9.3), Inches(0.5),
        "谢谢聆听", sz=28, c=C_WHITE, b=True, al=PP_ALIGN.CENTER)

    rect(s, Inches(0), Inches(7.0), SW, Inches(0.5), fill=C_PRI_D)


# ============================================================
# main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for fn in [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
               slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
               slide_13, slide_14, slide_15, slide_16, slide_17]:
        fn(prs)

    base = r"D:\Project\git"
    v = 10
    while True:
        out = os.path.join(base, f"竞聘汇报PPT_v{v}.pptx")
        if not os.path.exists(out):
            break
        v += 1
    prs.save(out)
    print(f"PPT generated: {out}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
